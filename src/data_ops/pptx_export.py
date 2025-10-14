"""
PowerPoint Export
Export analysis results to PowerPoint presentations

SEPARATION OF CONCERNS: PowerPoint generation only
Note: Requires python-pptx package
"""

import pandas as pd
from datetime import datetime
import io


class PowerPointExporter:
    """Export data and charts to PowerPoint"""
    
    @staticmethod
    def create_presentation(df, file_path, title="Data Analysis Report"):
        """
        Create a basic PowerPoint presentation
        
        Parameters:
        -----------
        df : DataFrame
            Data to include
        file_path : str
            Output file path
        title : str
            Presentation title
        
        Returns:
        --------
        success : bool
        message : str
        """
        try:
            # Try to import python-pptx
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN
            except ImportError:
                return False, "python-pptx package not installed. Install with: pip install python-pptx"
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title_shape.text = title
            subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}\nNexData Analytics"
            
            # Summary slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = "Data Summary"
            
            tf = body_shape.text_frame
            tf.text = f"Total Records: {len(df):,}"
            
            p = tf.add_paragraph()
            p.text = f"Columns: {len(df.columns)}"
            
            p = tf.add_paragraph()
            p.text = f"Date Range: {df.select_dtypes(include=['datetime64']).iloc[:, 0].min() if len(df.select_dtypes(include=['datetime64']).columns) > 0 else 'N/A'} to {df.select_dtypes(include=['datetime64']).iloc[:, 0].max() if len(df.select_dtypes(include=['datetime64']).columns) > 0 else 'N/A'}"
            
            # Statistics slide
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]
                
                title_shape.text = "Key Statistics"
                
                tf = body_shape.text_frame
                stats = df[numeric_cols].describe()
                
                for col in numeric_cols[:5]:  # First 5 numeric columns
                    p = tf.add_paragraph() if tf.text else tf
                    if tf.text:
                        p.level = 0
                    p.text = f"{col}:"
                    
                    sub_p = tf.add_paragraph()
                    sub_p.level = 1
                    sub_p.text = f"Mean: {stats.loc['mean', col]:.2f}, Median: {stats.loc['50%', col]:.2f}"
            
            # Save presentation
            prs.save(file_path)
            
            return True, f"Presentation created: {file_path}"
        
        except Exception as e:
            return False, f"Error creating presentation: {str(e)}"
    
    @staticmethod
    def add_data_table_slide(prs, title, df, max_rows=10):
        """
        Add a slide with data table
        
        Note: This is a helper method for future enhancements
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Add table
            rows, cols = min(len(df), max_rows) + 1, len(df.columns)
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(0.5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Header row
            for col_idx, col_name in enumerate(df.columns):
                table.cell(0, col_idx).text = str(col_name)
            
            # Data rows
            for row_idx in range(min(len(df), max_rows)):
                for col_idx, col_name in enumerate(df.columns):
                    table.cell(row_idx + 1, col_idx).text = str(df.iloc[row_idx, col_idx])
            
            return True
        
        except Exception as e:
            return False
    
    @staticmethod
    def check_pptx_available():
        """Check if python-pptx is installed"""
        try:
            import pptx
            return True, "python-pptx is installed"
        except ImportError:
            return False, "python-pptx not installed. Install with: pip install python-pptx"
